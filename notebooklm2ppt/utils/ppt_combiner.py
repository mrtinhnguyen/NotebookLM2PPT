# Sử dụng Spire.Presentation để hợp nhất các tệp PPT, giữ nguyên thiết kế gốc
# Cài đặt: pip install spire.presentation

import os
from pathlib import Path
from spire.presentation import *
from spire.presentation.common import *
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation as PptxPresentation
from PIL import Image

def combine_ppt_files_with_spire(source_folder, output_file, png_names=None):
    """
    Dùng Spire.Presentation để hợp nhất các file PPT, mỗi PPT chỉ giữ trang đầu và giữ nguyên thiết kế gốc.

    Args:
        source_folder: Thư mục chứa các file PPT nguồn
        output_file: Đường dẫn file PPT hợp nhất đầu ra
    """
    # Lấy tất cả các file .pptx và sắp xếp theo thứ tự từ điển
    
    ppt_files = sorted([f for f in os.listdir(source_folder) if f.endswith('.pptx')])

    ppt_names = [png_name.replace(".png", ".pptx") for png_name in png_names] if png_names else None
    
    if not ppt_files:
        print("Không tìm thấy file PPT nào")
        return
    
    print(f"Tìm thấy {len(ppt_files)} file PPT:")
    valid_ppt_files = []
    for idx, file in enumerate(ppt_files, 1):
        print(f"  {idx}. {file}")
        if ppt_names and file in ppt_names:
            valid_ppt_files.append(file)
    if ppt_names:
        ppt_files = valid_ppt_files
        print(f"\nSau khi lọc theo tên PNG được cung cấp, còn {len(ppt_files)} file PPT:")
    
    # Tạo đối tượng trình chiếu chính, dùng PPT đầu tiên làm cơ sở
    first_ppt_path = os.path.join(source_folder, ppt_files[0])
    main_pres = Presentation()
    main_pres.LoadFromFile(first_ppt_path)
    
    # Xóa các slide thừa trong PPT đầu tiên, chỉ giữ slide đầu
    while main_pres.Slides.Count > 1:
        main_pres.Slides.RemoveAt(1)
    
    print(f"  Đã thêm: {ppt_files[0]} (trang 1)")
    
    # Duyệt các file PPT còn lại
    for ppt_file in ppt_files[1:]:
        file_path = os.path.join(source_folder, ppt_file)
        
        # Tải PPT hiện tại
        temp_pres = Presentation()
        temp_pres.LoadFromFile(file_path)
        
        if temp_pres.Slides.Count > 0:
            # Dùng AppendBySlide để thêm trang đầu, giữ nguyên thiết kế
            main_pres.Slides.AppendBySlide(temp_pres.Slides[0])
            print(f"  Đã thêm: {ppt_file} (trang 1)")
        else:
            print(f"  Bỏ qua: {ppt_file} (không có slide)")
        
        # Giải phóng tài nguyên trình chiếu tạm thời
        temp_pres.Dispose()
    
    # Lưu PPT đã hợp nhất
    main_pres.SaveToFile(output_file, FileFormat.Pptx2016)
    print(f"\nĐã hợp nhất! Tệp đầu ra: {output_file}")
    print(f"Tổng cộng đã hợp nhất {main_pres.Slides.Count} slide")
    
    # Giải phóng tài nguyên
    main_pres.Dispose()

    valid_png_names = [ppt_name.replace(".pptx", ".png") for ppt_name in ppt_files] # Trả về tên PNG thực tế tồn tại

    return valid_png_names



def clean_ppt(in_ppt_file, out_ppt_file):
    """
    Dùng python-pptx để xóa các shape thừa trong PPT

    Args:
        ppt_file: Đường dẫn file PPT cần dọn dẹp
    """
    ppt = PptxPresentation(in_ppt_file)
    for slide in ppt.slides:
        for shape in list(slide.shapes):
            # Nếu tên shape là "New shape" thì xóa nó
            if shape.name == "New shape":
                sp = slide.shapes._spTree.remove(shape._element)

    ppt.save(out_ppt_file)
    

def combine_ppt(source_folder, out_ppt_file, png_names = None):
    # Đảm bảo là chuỗi đường dẫn vì sau này dùng .replace
    source_folder = str(source_folder)
    out_ppt_file = str(out_ppt_file)
    
    # Phương pháp 1: Giữ nguyên thiết kế gốc (khuyến nghị)
    output_file1 = out_ppt_file.replace(".pptx", "_combined_original_design.pptx")
    valid_png_names = combine_ppt_files_with_spire(source_folder, output_file1, png_names=png_names)

    clean_ppt(output_file1, out_ppt_file)
    print(f"\nĐã tạo file PPT hợp nhất: {out_ppt_file}")
    os.remove(output_file1)
    
    return valid_png_names


def create_ppt_from_images(png_dir, output_file, png_names=None):
    """
    Chèn trực tiếp ảnh PNG vào PPTX, không sử dụng nhận vùng thông minh

    Args:
        png_dir: Thư mục chứa ảnh PNG
        output_file: Đường dẫn file PPT đầu ra
        png_names: Danh sách tên file PNG (tuỳ chọn; nếu có thì chỉ xử lý những file này)

    Returns:
        Danh sách tên PNG thực tế đã sử dụng
    """
    png_dir = Path(png_dir)
    output_file = str(output_file)
    
    if png_names:
        png_files = [png_dir / name for name in png_names if (png_dir / name).exists()]
    else:
        png_files = sorted(png_dir.glob("*.png"))
    
    if not png_files:
        print("Không tìm thấy PNG nào")
        return []
    
    print(f"Tìm thấy {len(png_files)} ảnh PNG")
    
    prs = PptxPresentation()
    
    first_img = Image.open(png_files[0])
    img_width_px, img_height_px = first_img.size
    
    prs.slide_width = Pt(img_width_px)
    prs.slide_height = Pt(img_height_px)
    
    print(f"Kích thước PPT đặt là: {img_width_px} x {img_height_px} px")
    
    blank_layout = prs.slide_layouts[6]
    
    for idx, png_file in enumerate(png_files, 1):
        print(f"  [{idx}/{len(png_files)}] Xử lý: {png_file.name}")
        
        slide = prs.slides.add_slide(blank_layout)
        
        img = Image.open(png_file)
        img_width, img_height = img.size
        
        slide.shapes.add_picture(str(png_file), 0, 0, Pt(img_width), Pt(img_height))
    
    prs.save(output_file)
    print(f"\nĐã tạo file PPT: {output_file}")
    print(f"Tổng cộng đã thêm {len(prs.slides)} slide")
    
    return [f.name for f in png_files]
